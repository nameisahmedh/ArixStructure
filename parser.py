import pdfplumber
import fitz  # PyMuPDF
import io
from PIL import Image
import os
import re
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
import csv
import zipfile
import base64
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Create a directory to store extracted images
try:
    if not os.path.exists("temp_images"):
        os.makedirs("temp_images")
except OSError as e:
    logger.warning(f"Could not create temp_images directory: {e}")

def clean_text(text):
    """Cleans up extracted text."""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def parse_document(file_bytes, filename):
    """Main router function that parses a file from bytes and extracts content."""
    extension = os.path.splitext(filename)[1].lower()
    
    parsers = {
        '.pdf': _parse_pdf,
        '.docx': _parse_docx,
        '.pptx': _parse_pptx,
        '.txt': _parse_txt,
        '.html': _parse_html,
        '.htm': _parse_html,
        '.csv': _parse_csv
    }
    
    if extension in parsers:
        return parsers[extension](file_bytes)
    else:
        logger.error(f"Unsupported file type: {extension}")
        return {
            "full_text": f"Error: Unsupported file type '{extension}'. Please upload a supported file format.",
            "tables": [],
            "image_files": [],
            "metadata": {"error": f"Unsupported format: {extension}"}
        }

def _parse_pdf(file_bytes):
    """Enhanced PDF parser with better data structuring."""
    all_text = ""
    all_tables = []
    image_files = []
    metadata = {"pages": 0, "extraction_method": "pdfplumber + fitz", "images_found": 0}
    
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            metadata["pages"] = len(pdf.pages)
            logger.info(f"Parsing {len(pdf.pages)} pages...")
            
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    all_text += f"--- PAGE {i+1} ---\n"
                    all_text += clean_text(page_text)
                    all_text += f"\n--- END PAGE {i+1} ---\n\n"
                
                tables = page.extract_tables()
                for table in tables:
                    if table:
                        clean_table = []
                        for row in table:
                            clean_row = [str(cell).strip() if cell is not None else "" for cell in row]
                            clean_table.append(clean_row)
                        all_tables.append(clean_table)

    except (FileNotFoundError, PermissionError) as e:
        logger.error(f"File access error with pdfplumber: {e}")
        metadata["extraction_error"] = str(e)
    except Exception as e:
        logger.error(f"Parsing error with pdfplumber: {e}")
        metadata["extraction_error"] = str(e)

    # Extract images with PyMuPDF
    image_filenames = set()
    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            for page_index in range(len(doc)):
                image_list = doc.get_page_images(page_index)
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # Sanitize filename to prevent path traversal
                        safe_ext = re.sub(r'[^a-zA-Z0-9]', '', str(image_ext))
                        if not safe_ext:
                            safe_ext = 'png'
                        
                        # Secure filename generation
                        safe_filename = f"pdf_p{page_index+1}_{img_index}.{safe_ext}"
                        image_filename = os.path.join("temp_images", os.path.basename(safe_filename))
                        
                        # Validate and save image
                        if _validate_and_save_image(image_bytes, image_filename):
                            image_filenames.add(image_filename)
                            
                    except (IOError, OSError) as e:
                        logger.error(f"Image file error: {e}")
                    except Exception as e:
                        logger.error(f"Image extraction error: {e}")
    except Exception as e:
        logger.error(f"Error with image extraction: {e}")

    image_files = sorted(list(image_filenames))
    metadata["images_found"] = len(image_files)
    
    return {
        "full_text": all_text,
        "tables": all_tables,
        "image_files": image_files,
        "metadata": metadata
    }

def _parse_docx(file_bytes):
    """Enhanced DOCX parser with image extraction and table detection."""
    all_text = ""
    all_tables = []
    image_files = []
    metadata = {"extraction_method": "python-docx", "tables_found": 0}
    
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        
        # Extract images from DOCX
        image_files = _extract_docx_images(file_bytes)
        
        # Extract real tables first
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                all_tables.append(table_data)
        
        # Extract text and detect CSV-like structures
        current_csv_table = []
        for para in doc.paragraphs:
            text = para.text.strip()
            all_text += text + "\n"

            if not text:
                continue
            
            # Detect structured data patterns
            if _is_structured_data(text):
                delimiter = ',' if ',' in text else '\t'
                try:
                    f = io.StringIO(text)
                    reader = csv.reader(f, delimiter=delimiter)
                    for row in reader:
                        cleaned_row = [cell.strip() for cell in row if cell.strip()]
                        if len(cleaned_row) > 1:
                            current_csv_table.append(cleaned_row)
                except (ValueError, csv.Error):
                    pass
            else:
                if current_csv_table:
                    all_tables.append(current_csv_table)
                    current_csv_table = []
        
        # Add final table if exists
        if current_csv_table:
            all_tables.append(current_csv_table)
        
        metadata["tables_found"] = len(all_tables)
        metadata["images_found"] = len(image_files)
            
    except Exception as e:
        logger.error(f"Error parsing DOCX: {e}")
        metadata["extraction_error"] = str(e)
    
    return {
        "full_text": clean_text(all_text),
        "tables": all_tables,
        "image_files": image_files,
        "metadata": metadata
    }

def _parse_pptx(file_bytes):
    """Enhanced PPTX parser with image extraction and slide structure preservation."""
    all_text = ""
    all_tables = []
    image_files = []
    metadata = {"extraction_method": "python-pptx", "slides": 0}
    
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        metadata["slides"] = len(prs.slides)
        
        # Extract images from PPTX
        image_files = _extract_pptx_images(file_bytes)
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = f"--- SLIDE {slide_idx + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    slide_text += shape.text + "\n"
                
                if hasattr(shape, "table") and shape.table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data:
                        all_tables.append(table_data)
            
            slide_text += f"--- END SLIDE {slide_idx + 1} ---\n\n"
            all_text += slide_text
        
        metadata["images_found"] = len(image_files)
            
    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        metadata["extraction_error"] = str(e)
    
    return {
        "full_text": clean_text(all_text),
        "tables": all_tables,
        "image_files": image_files,
        "metadata": metadata
    }

def _parse_txt(file_bytes):
    """Enhanced TXT parser with automatic structure detection."""
    all_text = ""
    all_tables = []
    metadata = {"extraction_method": "text_analysis"}
    
    try:
        all_text = file_bytes.decode('utf-8', errors='ignore')
        
        # Detect structured data in text
        current_table = []
        lines = all_text.splitlines()
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if _is_structured_data(line):
                delimiter = ',' if line.count(',') >= 2 else '\t'
                try:
                    f = io.StringIO(line)
                    reader = csv.reader(f, delimiter=delimiter)
                    for row in reader:
                        cleaned_row = [cell.strip() for cell in row if cell.strip()]
                        if len(cleaned_row) > 1:
                            current_table.append(cleaned_row)
                except (ValueError, csv.Error):
                    pass
            else:
                if current_table:
                    all_tables.append(current_table)
                    current_table = []
        
        if current_table:
            all_tables.append(current_table)
            
    except Exception as e:
        logger.error(f"Error parsing TXT: {e}")
        metadata["extraction_error"] = str(e)
    
    return {
        "full_text": clean_text(all_text),
        "tables": all_tables,
        "image_files": [],
        "metadata": metadata
    }

def _parse_html(file_bytes):
    """Enhanced HTML parser with table and image extraction."""
    all_text = ""
    all_tables = []
    image_files = []
    metadata = {"extraction_method": "beautifulsoup"}
    
    try:
        html_content = file_bytes.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Extract images from HTML
        image_files = _extract_html_images(soup)
        
        # Extract tables
        tables = soup.find_all('table')
        for table in tables:
            table_data = []
            rows = table.find_all('tr')
            for row in rows:
                cells = row.find_all(['td', 'th'])
                row_data = [cell.get_text().strip() for cell in cells]
                if row_data:
                    table_data.append(row_data)
            
            if table_data:
                all_tables.append(table_data)
        
        # Extract text content
        all_text = soup.get_text()
        
        metadata["tables_found"] = len(all_tables)
        metadata["images_found"] = len(image_files)
        
    except Exception as e:
        logger.error(f"Error parsing HTML: {e}")
        metadata["extraction_error"] = str(e)
    
    return {
        "full_text": clean_text(all_text),
        "tables": all_tables,
        "image_files": image_files,
        "metadata": metadata
    }

def _parse_csv(file_bytes):
    """Enhanced CSV parser."""
    all_text = ""
    all_tables = []
    metadata = {"extraction_method": "csv_reader"}
    
    try:
        csv_content = file_bytes.decode('utf-8', errors='ignore')
        all_text = csv_content
        
        # Parse CSV data
        f = io.StringIO(csv_content)
        reader = csv.reader(f)
        table_data = []
        
        for row in reader:
            cleaned_row = [cell.strip() for cell in row]
            table_data.append(cleaned_row)
        
        if table_data:
            all_tables.append(table_data)
        
        metadata["tables_found"] = len(all_tables)
        
    except Exception as e:
        logger.error(f"Error parsing CSV: {e}")
        metadata["extraction_error"] = str(e)
    
    return {
        "full_text": clean_text(all_text),
        "tables": all_tables,
        "image_files": [],
        "metadata": metadata
    }

def _is_structured_data(text):
    """Detect if text contains structured data patterns."""
    if not text or len(text) < 10:
        return False
    
    # Check for common delimiters
    delimiters = [',', '\t', '|', ';']
    for delimiter in delimiters:
        if text.count(delimiter) >= 2:
            return True
    
    # Check for key-value patterns
    if ':' in text and ('=' in text or text.count(':') >= 2):
        return True
    
    return False

def _extract_docx_images(file_bytes):
    """Extract images from DOCX file."""
    image_files = []
    
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes), 'r') as docx_zip:
            # Look for images in the media folder
            for file_info in docx_zip.filelist:
                if file_info.filename.startswith('word/media/'):
                    # Extract image
                    image_data = docx_zip.read(file_info.filename)
                    
                    # Get file extension
                    _, ext = os.path.splitext(file_info.filename)
                    if not ext:
                        ext = '.png'  # Default extension
                    
                    # Create safe filename
                    safe_filename = f"docx_img_{len(image_files)}{ext}"
                    image_path = os.path.join("temp_images", os.path.basename(safe_filename))
                    
                    # Validate and save image
                    if _validate_and_save_image(image_data, image_path):
                        image_files.append(image_path)
                        
    except Exception as e:
        logger.error(f"Error extracting DOCX images: {e}")
    
    return image_files

def _extract_pptx_images(file_bytes):
    """Extract images from PPTX file."""
    image_files = []
    
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes), 'r') as pptx_zip:
            # Look for images in the media folder
            for file_info in pptx_zip.filelist:
                if file_info.filename.startswith('ppt/media/'):
                    # Extract image
                    image_data = pptx_zip.read(file_info.filename)
                    
                    # Get file extension
                    _, ext = os.path.splitext(file_info.filename)
                    if not ext:
                        ext = '.png'  # Default extension
                    
                    # Create safe filename
                    safe_filename = f"pptx_img_{len(image_files)}{ext}"
                    image_path = os.path.join("temp_images", os.path.basename(safe_filename))
                    
                    # Validate and save image
                    if _validate_and_save_image(image_data, image_path):
                        image_files.append(image_path)
                        
    except Exception as e:
        logger.error(f"Error extracting PPTX images: {e}")
    
    return image_files

def _extract_html_images(soup):
    """Extract images from HTML content."""
    image_files = []
    
    try:
        # Find all img tags
        img_tags = soup.find_all('img')
        
        for i, img in enumerate(img_tags):
            src = img.get('src', '')
            
            # Handle base64 encoded images
            if src.startswith('data:image/'):
                try:
                    # Extract base64 data
                    header, data = src.split(',', 1)
                    image_data = base64.b64decode(data)
                    
                    # Determine file extension from header
                    if 'jpeg' in header or 'jpg' in header:
                        ext = '.jpg'
                    elif 'png' in header:
                        ext = '.png'
                    elif 'gif' in header:
                        ext = '.gif'
                    else:
                        ext = '.png'  # Default
                    
                    # Create safe filename
                    safe_filename = f"html_img_{i}{ext}"
                    image_path = os.path.join("temp_images", os.path.basename(safe_filename))
                    
                    # Validate and save image
                    if _validate_and_save_image(image_data, image_path):
                        image_files.append(image_path)
                        
                except Exception as e:
                    logger.error(f"Error processing base64 image {i}: {e}")
                    
    except Exception as e:
        logger.error(f"Error extracting HTML images: {e}")
    
    return image_files

def _validate_and_save_image(image_data, image_path):
    """Validate image data and save if valid."""
    try:
        # Validate image data
        if len(image_data) < 100:  # Too small to be a valid image
            return False
        
        # Try to open with PIL to validate
        with Image.open(io.BytesIO(image_data)) as img:
            # Check image dimensions
            if img.width < 10 or img.height < 10:
                return False
            
            # Convert to RGB if necessary
            if img.mode in ('RGBA', 'LA', 'P'):
                img = img.convert('RGB')
            
            # Save the image
            img.save(image_path, format='PNG', optimize=True)
            return True
            
    except Exception as e:
        logger.error(f"Error validating/saving image: {e}")
        return False